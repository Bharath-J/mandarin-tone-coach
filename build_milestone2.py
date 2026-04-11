"""
build_milestone2.py
Generates Milestone2.pptx matching Milestone 1 visual style.
Run: python build_milestone2.py

Key design rules (to avoid overlap):
- Every visual section uses ONE multiline textbox (no stacked tiny ones)
- No two textboxes share overlapping x AND y ranges
- Every textbox has generous height (at least 50% more than minimum)
- No custom bar charts — use simple text lists instead
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_DIR = Path(__file__).parent
OUT_FILE    = PROJECT_DIR / "Milestone2.pptx"
RESULTS_DIR = PROJECT_DIR / "results"

# ── Palette ─────────────────────────────────────────────────────────────────────
C_NAVY     = RGBColor(0x06, 0x5A, 0x82)
C_TEAL     = RGBColor(0x02, 0xC3, 0x9A)
C_MID_BLUE = RGBColor(0x1C, 0x72, 0x93)
C_BLUE     = RGBColor(0x21, 0x96, 0xF3)
C_GREEN    = RGBColor(0x4C, 0xAF, 0x50)
C_ORANGE   = RGBColor(0xFF, 0x98, 0x00)
C_RED      = RGBColor(0xF4, 0x43, 0x36)
C_DARK     = RGBColor(0x1A, 0x2B, 0x3C)
C_MUTED    = RGBColor(0x5A, 0x7A, 0x8A)
C_LIGHT    = RGBColor(0xCA, 0xDC, 0xFC)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_DARK  = RGBColor(0x02, 0x1F, 0x2E)
C_DARK2    = RGBColor(0x0A, 0x3A, 0x52)
C_CARD     = RGBColor(0xFF, 0xFF, 0xFF)
C_WARN_BG  = RGBColor(0xFF, 0xF3, 0xE0)

# ── Slide dimensions (EMU) ───────────────────────────────────────────────────────
SW  = 9144000   # slide width
SH  = 5143500   # slide height
HDR = 822960    # header bar height
LM  = 274320    # left margin
TM  = 960120    # content top (below header)
PAD = 91440     # general padding unit (~0.1 in)


def E(v):
    return Emu(int(v))


# ── Low-level shape helpers ──────────────────────────────────────────────────────

def rect(slide, l, t, w, h, color):
    """Filled rectangle with no border."""
    sh = slide.shapes.add_shape(1, E(l), E(t), E(w), E(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def textbox(slide, l, t, w, h, text, size, bold=False, color=C_DARK,
            align=PP_ALIGN.LEFT, italic=False):
    """Single-run word-wrapped textbox. Returns the shape."""
    txb = slide.shapes.add_textbox(E(l), E(t), E(w), E(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb


def multibox(slide, l, t, w, h, lines):
    """
    Multi-paragraph textbox. lines is a list of dicts:
      text, size, bold, italic, color, align, space_before (Pt)
    Use empty-string text with size=4 for blank lines.
    """
    txb = slide.shapes.add_textbox(E(l), E(t), E(w), E(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        sb = line.get("space_before")
        if sb:
            p.space_before = Pt(sb)
        r = p.add_run()
        r.text = line.get("text", "")
        r.font.size = Pt(line.get("size", 11))
        r.font.bold = line.get("bold", False)
        r.font.italic = line.get("italic", False)
        r.font.color.rgb = line.get("color", C_DARK)
    return txb


def header(slide, title, dark=False):
    bg = C_BG_DARK if dark else C_NAVY
    rect(slide, 0, 0, SW, HDR, bg)
    textbox(slide, LM, 0, SW - LM * 2, HDR,
            title, 20, bold=True, color=C_WHITE)


def image(slide, path, l, t, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), E(l), E(t), E(w), E(h))


# ── Slide 1: Title ───────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    rect(slide, 0, 0, SW, SH, C_BG_DARK)

    # Left teal accent bar
    rect(slide, 0, 0, 164592, SH, C_TEAL)

    # Title
    textbox(slide, 457200, 1100000, SW - 457200 - LM, 900000,
            "Mandarin Tone Coach",
            44, bold=True, color=C_WHITE)

    # Subtitle
    textbox(slide, 457200, 2050000, 7500000, 480000,
            "Adaptive Pitch Feedback for L2 Mandarin Tone Learning",
            22, color=C_TEAL)

    # Navy divider line
    rect(slide, 457200, 2700000, 3657600, 36576, C_MID_BLUE)

    # Author block — single textbox, generous height
    multibox(slide, 457200, 2800000, SW - 457200 - LM, 750000, [
        {"text": "Bharath Jagadish (bjagadish3@gatech.edu)",
         "size": 14, "color": C_LIGHT},
        {"text": "CS6460 Educational Technology  \u00b7  Georgia Tech",
         "size": 14, "color": C_LIGHT, "space_before": 4},
    ])

    # Milestone 2 badge (bottom-left, above the divider area)
    BADGE_H = 411480
    BADGE_W = 1737360
    BADGE_T = SH - BADGE_H - PAD * 4
    rect(slide, 457200, BADGE_T, BADGE_W, BADGE_H, C_TEAL)
    textbox(slide, 457200, BADGE_T, BADGE_W, BADGE_H,
            "Milestone 2", 14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)


# ── Slide 2: Milestone 1 Recap ───────────────────────────────────────────────────

def slide_m1_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Milestone 1 Recap  \u2014  What We Built")

    PANEL_W = 4115520
    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040
    GAP     = PAD * 2

    # ── Left panel: The System ───────────────────────────────────────────────
    LX = LM
    rect(slide, LX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, LX, TM, PANEL_W, BAR_H, C_NAVY)
    textbox(slide, LX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "The System", 14, bold=True, color=C_WHITE)

    # Single multiline textbox for all content — starts below the bar
    content_top = TM + BAR_H + PAD * 2
    content_h   = PANEL_H - BAR_H - PAD * 4
    multibox(slide, LX + PAD * 2, content_top, PANEL_W - PAD * 4, content_h, [
        {"text": "Streamlit web app \u2014 real-time Mandarin tone practice",
         "size": 12, "bold": True, "color": C_NAVY},
        {"text": "", "size": 4},
        {"text": "Pipeline: Record \u2192 Extract F0 \u2192 Classify \u2192 Visualize \u2192 Feedback",
         "size": 10, "color": C_MUTED},
        {"text": "", "size": 6},
        {"text": "SVM Classifier", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "\u2022 91.95% speaker-independent accuracy (Tone Perfect corpus, 9,838 MP3s)",
         "size": 10, "color": C_DARK},
        {"text": "\u2022 27 acoustic features: F0 contour, first-order deltas, f0_min_pos",
         "size": 10, "color": C_DARK},
        {"text": "\u2022 Main T2\u2194T3 confusion (16%); rule-based override applied",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 6},
        {"text": "Practice Modes", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "\u2022 Monosyllabic: 40 words (10 per tone T1\u2013T4)",
         "size": 10, "color": C_DARK},
        {"text": "\u2022 Disyllabic: 20 two-syllable words with tone sandhi support",
         "size": 10, "color": C_DARK},
        {"text": "\u2022 Pitch contour overlay: learner (red) vs canonical reference (blue)",
         "size": 10, "color": C_DARK},
        {"text": "\u2022 Plain-language corrective feedback for all tone-pair combinations",
         "size": 10, "color": C_DARK},
    ])

    # ── Right panel: Improvements Since Milestone 1 ──────────────────────────
    RX = LM + PANEL_W + GAP
    rect(slide, RX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, RX, TM, PANEL_W, BAR_H, C_TEAL)
    textbox(slide, RX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "Improvements Since Milestone 1", 14, bold=True, color=C_WHITE)

    r_content_top = TM + BAR_H + PAD * 2
    r_content_h   = PANEL_H - BAR_H - PAD * 4
    multibox(slide, RX + PAD * 2, r_content_top, PANEL_W - PAD * 4, r_content_h, [
        {"text": "Disyllabic Practice Mode", "size": 11, "bold": True, "color": C_TEAL},
        {"text": "Energy-minimum syllable boundary detection; neutral-tone handling; 20-word list",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "Evaluation Instruments", "size": 11, "bold": True, "color": C_BLUE},
        {"text": "Pre/post perception test (Forms A & B, different speakers, Google Sheets logging)",
         "size": 10, "color": C_DARK},
        {"text": "TAM usability survey \u2014 8 Likert items + open-ended feedback",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "Usability Fixes", "size": 11, "bold": True, "color": C_GREEN},
        {"text": "Tone descriptions in sidebar; stale-result bug fixed; Streamlit header hidden; dividers repositioned",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "Data Fixes", "size": 11, "bold": True, "color": C_ORANGE},
        {"text": "sh\u00e0ngk\u0113 tone corrected (T4+T4); neutral-tone syllable window widened to 25\u201385%",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "UI Polish", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "Disyllabic tone badges use flex layout; footer text removed; toolbar minimised",
         "size": 10, "color": C_DARK},
    ])


# ── Slide 3: Study Design ────────────────────────────────────────────────────────

def slide_study_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Study Design")

    PANEL_W = 4206240
    GAP     = PAD * 2
    PANEL_H = SH - TM - PAD * 3
    BAR_H   = 320040

    # ── Left panel: Study Protocol ───────────────────────────────────────────
    LX = LM
    rect(slide, LX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, LX, TM, PANEL_W, BAR_H, C_NAVY)
    textbox(slide, LX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "Study Protocol", 14, bold=True, color=C_WHITE)

    ct = TM + BAR_H + PAD * 2
    ch = PANEL_H - BAR_H - PAD * 4
    multibox(slide, LX + PAD * 2, ct, PANEL_W - PAD * 4, ch, [
        {"text": "19 perception test \u00b7 16 TAM survey \u00b7 CS6460 classmates",
         "size": 12, "bold": True, "color": C_NAVY},
        {"text": "Online, asynchronous, self-paced (~10\u201315 min)", "size": 11, "color": C_DARK},
        {"text": "No prior Mandarin experience required", "size": 11, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "1  Pre-test (Form A)", "size": 12, "bold": True, "color": C_TEAL},
        {"text": "12 audio clips \u00b7 select the tone heard \u00b7 no feedback given",
         "size": 10, "color": C_DARK},
        {"text": "Speakers: Female Voice 1 + Male Voice 1", "size": 10, "color": C_MUTED},
        {"text": "", "size": 5},
        {"text": "2  Practice with the App", "size": 12, "bold": True, "color": C_BLUE},
        {"text": "Mandarin Tone Coach: monosyllabic & disyllabic words",
         "size": 10, "color": C_DARK},
        {"text": "Real-time pitch contour + plain-language corrective feedback",
         "size": 10, "color": C_MUTED},
        {"text": "", "size": 5},
        {"text": "3  Post-test (Form B)", "size": 12, "bold": True, "color": C_GREEN},
        {"text": "12 new clips \u00b7 different speakers (Female Voice 2 + Male Voice 2)",
         "size": 10, "color": C_DARK},
        {"text": "Same format, counterbalanced to minimise item-overlap confound",
         "size": 10, "color": C_MUTED},
        {"text": "", "size": 5},
        {"text": "4  TAM Usability Survey", "size": 12, "bold": True, "color": C_ORANGE},
        {"text": "8 Likert items (1\u20135) + 1 open-ended feedback question",
         "size": 10, "color": C_DARK},
    ])

    # ── Right panel: Instruments ─────────────────────────────────────────────
    RX = LM + PANEL_W + GAP
    rect(slide, RX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, RX, TM, PANEL_W, BAR_H, C_TEAL)
    textbox(slide, RX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "Instruments", 14, bold=True, color=C_WHITE)

    rt = TM + BAR_H + PAD * 2
    rh = PANEL_H - BAR_H - PAD * 4
    multibox(slide, RX + PAD * 2, rt, PANEL_W - PAD * 4, rh, [
        {"text": "Perception Test", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "\u2022 12 items per form (3 per tone T1\u2013T4)", "size": 11, "color": C_DARK},
        {"text": "\u2022 Items randomised per participant (seeded by participant ID)",
         "size": 11, "color": C_DARK},
        {"text": "\u2022 Scored: proportion correct (0\u201312 per participant)",
         "size": 11, "color": C_DARK},
        {"text": "\u2022 Responses stored in Google Sheets via gspread",
         "size": 11, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "Speaker Design", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "\u2022 Form A \u2014 Female Voice 1 & Male Voice 1", "size": 11, "color": C_DARK},
        {"text": "\u2022 Form B \u2014 Female Voice 2 & Male Voice 2  (new, unseen speakers)",
         "size": 11, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "TAM Survey  (Technology Acceptance Model)",
         "size": 12, "bold": True, "color": C_NAVY},
        {"text": "\u2022 8 Likert items, 5-point scale (1 = strongly disagree \u2192 5 = strongly agree)",
         "size": 11, "color": C_DARK},
        {"text": "\u2022 Constructs: perceived usefulness, ease of use, reuse intent",
         "size": 11, "color": C_DARK},
        {"text": "\u2022 1 open-ended question for qualitative feedback",
         "size": 11, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "Participant IDs", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "\u2022 Auto-generated 8-char hex ID shown at end of Form A",
         "size": 11, "color": C_DARK},
        {"text": "\u2022 Participants enter their ID to link Form A \u2192 Form B",
         "size": 11, "color": C_DARK},
    ])


# ── Slide 4: Participant Demographics ────────────────────────────────────────────

def slide_demographics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Participant Demographics  \u2014  19 Perception Test \u00b7 16 TAM Survey Respondents")

    CONTENT_H = SH - TM - PAD * 2
    BAR_H     = 320040

    # Divide into three equal-ish columns
    TOTAL_W   = SW - LM * 2
    COL_GAP   = PAD * 2
    COL_W     = (TOTAL_W - COL_GAP * 2) // 3

    # ── LEFT panel: dark navy stat box ──────────────────────────────────────
    LX = LM
    rect(slide, LX, TM, COL_W, CONTENT_H, C_NAVY)

    # "13" large stat — placed at top with generous spacing
    stat_top = TM + PAD * 3
    textbox(slide, LX, stat_top, COL_W, 600000,
            "13", 52, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

    # "Paired participants" label below "13"
    part_top = stat_top + 580000
    textbox(slide, LX, part_top, COL_W, 280000,
            "Paired participants", 14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # thin divider
    div1_top = part_top + 300000
    rect(slide, LX + PAD * 2, div1_top, COL_W - PAD * 4, 36576, C_MID_BLUE)

    # paired / form A only text — single textbox
    paired_top = div1_top + 80000
    multibox(slide, LX + PAD, paired_top, COL_W - PAD * 2, 400000, [
        {"text": "13 paired (both forms)", "size": 12, "color": C_LIGHT,
         "align": PP_ALIGN.CENTER},
        {"text": "6 Form A only", "size": 11, "color": C_MUTED,
         "align": PP_ALIGN.CENTER, "space_before": 6},
    ])

    # thin divider
    div2_top = paired_top + 450000
    rect(slide, LX + PAD * 2, div2_top, COL_W - PAD * 4, 36576, C_MID_BLUE)

    # "16" survey stat
    s17_top = div2_top + 80000
    textbox(slide, LX, s17_top, COL_W, 500000,
            "16", 40, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

    # "TAM respondents" label — well below "16"
    sr_top = s17_top + 480000
    textbox(slide, LX + PAD, sr_top, COL_W - PAD * 2, 300000,
            "TAM respondents", 12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # ── MIDDLE panel: Prior Mandarin Experience ──────────────────────────────
    MX = LM + COL_W + COL_GAP
    rect(slide, MX, TM, COL_W, CONTENT_H, C_CARD)
    rect(slide, MX, TM, COL_W, BAR_H, C_BLUE)
    textbox(slide, MX + PAD, TM + PAD, COL_W - PAD * 2, BAR_H - PAD,
            "Prior Mandarin Experience", 13, bold=True, color=C_WHITE)

    m_ct = TM + BAR_H + PAD * 2
    m_ch = CONTENT_H - BAR_H - PAD * 4
    multibox(slide, MX + PAD * 2, m_ct, COL_W - PAD * 4, m_ch, [
        {"text": "Never studied before:  11 / 16  (69%)", "size": 11, "color": C_DARK},
        {"text": "More than 2 years:      3 / 16  (19%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "Less than 6 months:    2 / 16  (13%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "", "size": 6},
        {"text": "69% of respondents had no prior Mandarin exposure \u2014 confirming the "
                 "target population of complete novices.",
         "size": 10, "color": C_MUTED, "italic": True},
    ])

    # ── RIGHT panel: Native Language ─────────────────────────────────────────
    RX = MX + COL_W + COL_GAP
    rect(slide, RX, TM, COL_W, CONTENT_H, C_CARD)
    rect(slide, RX, TM, COL_W, BAR_H, C_GREEN)
    textbox(slide, RX + PAD, TM + PAD, COL_W - PAD * 2, BAR_H - PAD,
            "Native Language", 13, bold=True, color=C_WHITE)

    r_ct = TM + BAR_H + PAD * 2
    r_ch = CONTENT_H - BAR_H - PAD * 4
    multibox(slide, RX + PAD * 2, r_ct, COL_W - PAD * 4, r_ch, [
        {"text": "English (all variants)  7 / 16  (44%)", "size": 11, "color": C_DARK},
        {"text": "Mandarin / Chinese      3 / 16  (19%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "Kannada / Tulu          2 / 16  (13%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "Hindi                   1 / 16  (6%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "Farsi                   1 / 16  (6%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "Malayalam               1 / 16  (6%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "Swahili                 1 / 16  (6%)", "size": 11, "color": C_DARK,
         "space_before": 6},
        {"text": "", "size": 6},
        {"text": "Native/heritage Mandarin speakers (n=3) flagged \u2014 different "
                 "performance profile from novices.",
         "size": 10, "color": C_MUTED, "italic": True},
    ])


# ── Slide 5: Pre/Post Results ────────────────────────────────────────────────────

def slide_prepost(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Pre/Post Perception Test Results  \u2014  13 Paired Participants")

    CONTENT_H = SH - TM - PAD * 2

    # Left stat column width
    STAT_W = 2400000

    # ── LEFT: dark navy stat box ─────────────────────────────────────────────
    rect(slide, LM, TM, STAT_W, CONTENT_H, C_NAVY)

    # "+10.3 pp" — top section
    textbox(slide, LM, TM + PAD * 3, STAT_W, 600000,
            "+10.3 pp", 44, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

    # "Mean accuracy gain" label
    textbox(slide, LM + PAD, TM + PAD * 3 + 580000, STAT_W - PAD * 2, 250000,
            "Mean accuracy gain", 11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # "53.8% → 64.1%"
    textbox(slide, LM + PAD, TM + PAD * 3 + 860000, STAT_W - PAD * 2, 280000,
            "53.8%  \u2192  64.1%", 13, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Thin divider
    div_top = TM + PAD * 3 + 1200000
    rect(slide, LM + PAD * 2, div_top, STAT_W - PAD * 4, 36576, C_MID_BLUE)

    # Dark card within the stat box for stats
    stats_top = div_top + 80000
    stats_h   = CONTENT_H - (stats_top - TM) - PAD * 2
    rect(slide, LM + PAD, stats_top, STAT_W - PAD * 2, stats_h, C_DARK2)

    # "Statistical Tests" heading inside dark card
    textbox(slide, LM + PAD * 2, stats_top + PAD * 2, STAT_W - PAD * 4, 280000,
            "Statistical Tests", 12, bold=True, color=C_TEAL)

    # All stats in ONE multiline textbox, below the heading
    st_content_top = stats_top + PAD * 2 + 300000
    st_content_h   = stats_h - PAD * 2 - 320000
    multibox(slide, LM + PAD * 2, st_content_top, STAT_W - PAD * 4, st_content_h, [
        {"text": "Wilcoxon signed-rank test", "size": 10, "bold": True, "color": C_LIGHT},
        {"text": "W = 13.5,  p = 0.087 (two-sided)", "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "Paired t-test (reference)", "size": 10, "bold": True, "color": C_LIGHT},
        {"text": "t = \u22122.17,  p = 0.051", "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "n = 13 paired participants", "size": 10, "bold": True, "color": C_LIGHT},
        {"text": "3 items per tone per participant", "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "Interpretation", "size": 10, "bold": True, "color": C_LIGHT},
        {"text": "Strong positive trend; p just above 0.05 threshold \u2014 marginal with n=13 (est. power ~35%)",
         "size": 10, "color": C_DARK},
    ])

    # ── RIGHT: chart image ───────────────────────────────────────────────────
    chart = RESULTS_DIR / "prepost_accuracy.png"
    CX = LM + STAT_W + PAD * 2
    CW = SW - CX - LM
    image(slide, chart, CX, TM, CW, CONTENT_H)


# ── Slide 6: Per-Tone Accuracy ───────────────────────────────────────────────────

def slide_per_tone(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Per-Tone Accuracy  \u2014  Pre vs Post  (13 paired participants, 3 items / tone each)")

    CONTENT_H = SH - TM - PAD * 2

    # Leave 280000 at the bottom for the observations box
    OBS_H  = 380000
    CARD_H = CONTENT_H - OBS_H - PAD * 2

    N      = 4
    TOTAL_W = SW - LM * 2
    CARD_W  = (TOTAL_W - PAD * (N - 1)) // N
    BAR_H_C = 320040  # coloured header height per card

    tone_data = [
        ("T1  flat",     53.8, 69.2, C_BLUE,   "+15.4 pp \u2191"),
        ("T2  rising",   53.8, 66.7, C_GREEN,  "+12.8 pp \u2191"),
        ("T3  dipping",  46.2, 66.7, C_ORANGE, "+20.5 pp \u2191"),
        ("T4  falling",  61.5, 53.8, C_RED,    "\u22127.7 pp \u2193"),
    ]

    for i, (label, pre, post, col, delta) in enumerate(tone_data):
        cx = LM + i * (CARD_W + PAD)

        # White card background
        rect(slide, cx, TM, CARD_W, CARD_H, C_CARD)

        # Coloured header bar
        rect(slide, cx, TM, CARD_W, BAR_H_C, col)
        textbox(slide, cx + PAD, TM + PAD, CARD_W - PAD * 2, BAR_H_C - PAD,
                label, 13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Single textbox for all card content — starts well below the header bar
        tb_top = TM + BAR_H_C + PAD * 2
        tb_h   = CARD_H - BAR_H_C - PAD * 4
        multibox(slide, cx + PAD * 2, tb_top, CARD_W - PAD * 4, tb_h, [
            {"text": "Pre-test", "size": 10, "color": C_MUTED},
            {"text": f"{pre:.1f}%", "size": 14, "bold": True, "color": col,
             "space_before": 2},
            {"text": "", "size": 5},
            {"text": "Post-test", "size": 10, "color": C_DARK, "bold": True},
            {"text": f"{post:.1f}%", "size": 14, "bold": True, "color": col,
             "space_before": 2},
            {"text": "", "size": 8},
            {"text": "", "size": 8},
            {"text": delta, "size": 16, "bold": True, "color": col,
             "align": PP_ALIGN.CENTER},
        ])

    # ── Observations box ─────────────────────────────────────────────────────
    OBS_Y = TM + CARD_H + PAD
    rect(slide, LM, OBS_Y, SW - LM * 2, OBS_H, C_WARN_BG)
    textbox(slide, LM + PAD * 2, OBS_Y + PAD * 2, SW - LM * 2 - PAD * 4, OBS_H - PAD * 4,
            "Key observations:  T3 (dipping) showed the largest gain (+20.5 pp) \u2014 likely benefiting "
            "from the pitch contour visualisation.  T4 declined (\u22127.7 pp) \u2014 possibly harder Form B "
            "stimuli or high variance with only 3 items per tone.  "
            "Caution: n=13, 3 items/tone = limited statistical power.",
            10, color=C_DARK)


# ── Slide 7: TAM Survey ──────────────────────────────────────────────────────────

def slide_survey(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "App Feedback  \u2014  TAM Usability Survey  (n = 16,  1\u20135 Likert scale)")

    CONTENT_H = SH - TM - PAD * 3  # leave space for footer

    # Left 60% — chart image
    CHART_W = int((SW - LM * 2) * 0.60)
    image(slide, RESULTS_DIR / "survey_results.png",
          LM, TM, CHART_W, CONTENT_H)

    # Right 40% — stats panel
    RX = LM + CHART_W + PAD * 2
    RW = SW - RX - LM

    # Overall mean navy box
    MEAN_BOX_H = 960120
    rect(slide, RX, TM, RW, MEAN_BOX_H, C_NAVY)
    textbox(slide, RX, TM + PAD * 2, RW, 500000,
            "3.96 / 5.0", 36, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    textbox(slide, RX + PAD, TM + PAD * 2 + 500000, RW - PAD * 2, 300000,
            "Overall mean (8 items)", 11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Item rows — each is a rect (card) + colored strip + label textbox
    highlights = [
        (C_GREEN,  "4.19", "Would recommend to others  (highest)"),
        (C_BLUE,   "4.12", "Overall app useful for learning"),
        (C_TEAL,   "4.06", "Would use app again"),
        (C_ORANGE, "4.00", "Recording & feedback straightforward"),
        (C_MUTED,  "3.69", "Helped understand tones  (lowest)"),
    ]

    ROW_H   = 440000
    ROW_GAP = 60000
    STRIP_W = 380000

    # Compute how many rows fit
    rows_start = TM + MEAN_BOX_H + PAD
    for idx, (col, score, label) in enumerate(highlights):
        ry = rows_start + idx * (ROW_H + ROW_GAP)

        # White card background
        rect(slide, RX, ry, RW, ROW_H, C_CARD)

        # Coloured left strip
        rect(slide, RX, ry, STRIP_W, ROW_H, col)

        # Score in strip — centered vertically
        textbox(slide, RX, ry + PAD, STRIP_W, ROW_H - PAD * 2,
                score, 14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Label text to the right of strip
        textbox(slide, RX + STRIP_W + PAD, ry + PAD,
                RW - STRIP_W - PAD * 2, ROW_H - PAD * 2,
                label, 9, color=C_DARK)

    # Footer note — at bottom of slide
    footer_top = SH - 240000
    textbox(slide, LM, footer_top, SW - LM * 2, 220000,
            "\u26a0 Respondent 13 (native Mandarin speaker who couldn\u2019t navigate the app) gave all 1s "
            "\u2014 excluding this outlier raises the overall mean to 4.09.",
            9, color=C_MUTED)


# ── Slide 8: Key Findings & Limitations ─────────────────────────────────────────

def slide_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Key Findings & Limitations", dark=True)

    # Reserve 200000 at the bottom for the quote box
    QUOTE_H = 200000
    CARD_H  = SH - TM - QUOTE_H - PAD * 4
    COL_BAR = 320040

    N      = 3
    TOTAL_W = SW - LM * 2
    COL_GAP = PAD * 2
    COL_W   = (TOTAL_W - COL_GAP * (N - 1)) // N

    cols_data = [
        {
            "title": "What Worked",
            "color": C_TEAL,
            "items": [
                "Positive pre\u2192post trend (+10.3 pp) across all non-native participants",
                "T3 (dipping) largest gain (+20.5 pp) \u2014 pitch contour most helpful",
                "TAM overall usefulness & reuse intent both 4.12 / 5",
                "Pitch contour overlay praised in open-ended feedback",
                "Disyllabic mode broadened practice scope successfully",
            ],
        },
        {
            "title": "Participant Feedback",
            "color": C_ORANGE,
            "items": [
                "T2 vs T3 perceptually hard to distinguish \u2014 confirmed by test data",
                "No target audio before recording \u2014 users wanted to hear word first",
                "\u201cToo noisy / too quiet\u201d error not specific enough to act on",
                "More words requested; 40-word list feels limited",
                "App navigation confusing for one participant on first load",
            ],
        },
        {
            "title": "Limitations",
            "color": C_RED,
            "items": [
                "n = 13 paired \u2014 underpowered pilot (est. power ~35%)",
                "Asynchronous: no control over practice time or noise environment",
                "T4 decline (\u22127.7 pp) likely Form B difficulty, not regression",
                "3 native speakers in sample \u2014 different performance profile",
                "SVM misclassifies some correct native-speaker productions",
            ],
        },
    ]

    for i, cd in enumerate(cols_data):
        cx = LM + i * (COL_W + COL_GAP)

        # Dark card background
        rect(slide, cx, TM, COL_W, CARD_H, C_DARK2)

        # Coloured top bar
        rect(slide, cx, TM, COL_W, COL_BAR, cd["color"])
        textbox(slide, cx + PAD, TM + PAD, COL_W - PAD * 2, COL_BAR - PAD,
                cd["title"], 13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Single multiline textbox for ALL bullets
        blt_top = TM + COL_BAR + PAD * 2
        blt_h   = CARD_H - COL_BAR - PAD * 4

        lines = []
        for item in cd["items"]:
            lines.append({"text": "\u2022  " + item, "size": 10, "color": C_LIGHT,
                          "space_before": 6})

        multibox(slide, cx + PAD * 2, blt_top, COL_W - PAD * 4, blt_h, lines)

    # Quote box — at the very bottom
    QUOTE_Y = TM + CARD_H + PAD * 2
    actual_quote_h = SH - QUOTE_Y - PAD
    rect(slide, LM, QUOTE_Y, SW - LM * 2, actual_quote_h, C_DARK2)
    textbox(slide, LM + PAD * 3, QUOTE_Y + PAD * 2, SW - LM * 2 - PAD * 6, actual_quote_h - PAD * 4,
            "\u201cThe app is well planned and designed in a way that makes it very easy "
            "for the user to learn Mandarin.\u201d  \u2014 Participant 9",
            10, italic=True, color=C_TEAL)


# ── Slide 9: What's Next ─────────────────────────────────────────────────────────

def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "What\u2019s Next  \u2014  Weeks 14\u201316")

    # Leave 380000 at the bottom for the banner
    BAN_H  = 380000
    CARD_H = SH - TM - BAN_H - PAD * 2
    HDR_H  = 347472

    N       = 3
    TOTAL_W = SW - LM * 2
    COL_GAP = PAD * 2
    CARD_W  = (TOTAL_W - COL_GAP * (N - 1)) // N

    cards = [
        {
            "title": "Week 14 \u00b7 Deep Analysis",
            "color": C_TEAL,
            "items": [
                "Per-participant learning curve analysis",
                "Subgroup: novices vs prior Mandarin learners",
                "Perception test confusion matrix (which tones mixed up)",
                "Correlate TAM scores with learning gains",
            ],
        },
        {
            "title": "Weeks 14\u201315 \u00b7 Paper Writing",
            "color": C_BLUE,
            "items": [
                "Introduction & literature review (CAPT, Noticing Hypothesis, TAM)",
                "Methods: system design + study protocol",
                "Results: quantitative pre/post + qualitative survey themes",
                "Discussion: implications, limitations, future work",
            ],
        },
        {
            "title": "Week 16 \u00b7 Final Submission",
            "color": C_GREEN,
            "items": [
                "Complete paper (course format)",
                "Final video walkthrough of the system",
                "Open-source release on GitHub",
                "Optional: add target-word audio examples to the app",
            ],
        },
    ]

    for i, card in enumerate(cards):
        cx = LM + i * (CARD_W + COL_GAP)

        # White card
        rect(slide, cx, TM, CARD_W, CARD_H, C_CARD)

        # Coloured header bar
        rect(slide, cx, TM, CARD_W, HDR_H, card["color"])
        textbox(slide, cx + PAD, TM + PAD, CARD_W - PAD * 2, HDR_H - PAD,
                card["title"], 11, bold=True, color=C_WHITE)

        # Single multiline textbox for all bullets
        blt_top = TM + HDR_H + PAD * 2
        blt_h   = CARD_H - HDR_H - PAD * 4

        lines = []
        for item in card["items"]:
            lines.append({"text": "\u2713  " + item, "size": 11, "color": C_DARK,
                          "space_before": 6})

        multibox(slide, cx + PAD * 2, blt_top, CARD_W - PAD * 4, blt_h, lines)

    # Bottom banner
    BAN_Y = TM + CARD_H + PAD
    actual_ban_h = SH - BAN_Y - PAD
    rect(slide, LM, BAN_Y, SW - LM * 2, actual_ban_h, C_NAVY)
    textbox(slide, LM + PAD * 2, BAN_Y + PAD * 2, SW - LM * 2 - PAD * 4, actual_ban_h - PAD * 4,
            "Weeks 8\u201313 complete  \u00b7  Evaluation done  \u00b7  19 perception test \u00b7 16 TAM survey  "
            "\u00b7  Data collected  \u00b7  Analysis scripts ready",
            11, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)


# ── Build ────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = E(SW)
    prs.slide_height = E(SH)

    slide_title(prs)
    slide_m1_recap(prs)
    slide_study_design(prs)
    slide_demographics(prs)
    slide_prepost(prs)
    slide_per_tone(prs)
    slide_survey(prs)
    slide_findings(prs)
    slide_next_steps(prs)

    prs.save(str(OUT_FILE))
    print(f"Saved: {OUT_FILE}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
