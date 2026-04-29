"""
build_final_presentation.py
Generates Final_Presentation_bjagadish3.pptx — CS6460 final video presentation.
9 slides, ~10 minutes. Goal: make listener want to read the paper.

Design rules (same as build_milestone2.py):
- Every visual section uses ONE multiline textbox (no stacked tiny ones)
- No two textboxes share overlapping x AND y ranges
- Every textbox has generous height (at least 50% more than minimum)
- No custom bar charts — use simple text lists instead
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_DIR = Path(__file__).parent
OUT_FILE    = PROJECT_DIR / "Final_Presentation_bjagadish3.pptx"
RESULTS_DIR = PROJECT_DIR / "results"

# ── Palette ─────────────────────────────────────────────────────────────────────
C_NAVY     = RGBColor(0x06, 0x5A, 0x82)
C_TEAL     = RGBColor(0x02, 0xC3, 0x9A)
C_MID_BLUE = RGBColor(0x1C, 0x72, 0x93)
C_BLUE     = RGBColor(0x21, 0x96, 0xF3)
C_GREEN    = RGBColor(0x4C, 0xAF, 0x50)
C_ORANGE   = RGBColor(0xFF, 0x98, 0x00)
C_RED      = RGBColor(0xF4, 0x43, 0x36)
C_DARK     = RGBColor(0x1A, 0x2B, 0x3C)
C_MUTED    = RGBColor(0x5A, 0x7A, 0x8A)
C_LIGHT    = RGBColor(0xCA, 0xDC, 0xFC)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_DARK  = RGBColor(0x02, 0x1F, 0x2E)
C_DARK2    = RGBColor(0x0A, 0x3A, 0x52)
C_CARD     = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions (EMU) ───────────────────────────────────────────────────────
SW  = 9144000
SH  = 5143500
HDR = 822960
LM  = 274320
TM  = 960120
PAD = 91440


def E(v):
    return Emu(int(v))


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, E(l), E(t), E(w), E(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def textbox(slide, l, t, w, h, text, size, bold=False, color=C_DARK,
            align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(E(l), E(t), E(w), E(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb


def multibox(slide, l, t, w, h, lines):
    txb = slide.shapes.add_textbox(E(l), E(t), E(w), E(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        sb = line.get("space_before")
        if sb:
            p.space_before = Pt(sb)
        r = p.add_run()
        r.text = line.get("text", "")
        r.font.size = Pt(line.get("size", 11))
        r.font.bold = line.get("bold", False)
        r.font.italic = line.get("italic", False)
        r.font.color.rgb = line.get("color", C_DARK)
    return txb


def header(slide, title):
    rect(slide, 0, 0, SW, HDR, C_NAVY)
    textbox(slide, LM, 0, SW - LM * 2, HDR,
            title, 20, bold=True, color=C_WHITE)


def card(slide, lx, color=C_CARD, bar_color=C_NAVY,
         panel_w=4115520, bar_label="", bar_h=320040):
    panel_h = SH - TM - PAD * 2
    rect(slide, lx, TM, panel_w, panel_h, color)
    rect(slide, lx, TM, panel_w, bar_h, bar_color)
    if bar_label:
        textbox(slide, lx + PAD, TM + PAD, panel_w - PAD * 2, bar_h - PAD,
                bar_label, 14, bold=True, color=C_WHITE)
    return panel_h, bar_h


def image(slide, path, l, t, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), E(l), E(t), E(w), E(h))


# ── Slide 1: Title ───────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SW, SH, C_BG_DARK)
    rect(slide, 0, 0, 164592, SH, C_TEAL)

    textbox(slide, 457200, 900000, SW - 457200 - LM, 900000,
            "Mandarin Tone Coach", 44, bold=True, color=C_WHITE)
    textbox(slide, 457200, 1880000, 7500000, 480000,
            "Adaptive Pitch Feedback for L2 Mandarin Tone Learning",
            22, color=C_TEAL)
    rect(slide, 457200, 2550000, 3657600, 36576, C_MID_BLUE)
    multibox(slide, 457200, 2650000, SW - 457200 - LM, 700000, [
        {"text": "Bharath Jagadish  ·  bjagadish3@gatech.edu",
         "size": 14, "color": C_LIGHT},
        {"text": "CS6460 Educational Technology  ·  Georgia Tech",
         "size": 14, "color": C_LIGHT, "space_before": 4},
    ])

    BADGE_H = 411480
    BADGE_W = 2000000
    BADGE_T = SH - BADGE_H - PAD * 4
    rect(slide, 457200, BADGE_T, BADGE_W, BADGE_H, C_TEAL)
    textbox(slide, 457200, BADGE_T, BADGE_W, BADGE_H,
            "Final Presentation", 14, bold=True, color=C_NAVY,
            align=PP_ALIGN.CENTER)


# ── Slide 2: Problem Statement ───────────────────────────────────────────────────
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "The Problem — Why Mandarin Tones Are Hard")

    PANEL_W = 4115520
    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040
    GAP     = PAD * 2

    # Left card — The Challenge
    LX = LM
    rect(slide, LX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, LX, TM, PANEL_W, BAR_H, C_NAVY)
    textbox(slide, LX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "The Challenge", 14, bold=True, color=C_WHITE)

    ct = TM + BAR_H + PAD * 2
    ch = PANEL_H - BAR_H - PAD * 4
    multibox(slide, LX + PAD * 2, ct, PANEL_W - PAD * 4, ch, [
        {"text": "Mandarin is the world's most spoken language.",
         "size": 12, "bold": True, "color": C_NAVY},
        {"text": "", "size": 4},
        {"text": "Its four lexical tones make it notoriously difficult for "
                 "speakers of non-tonal languages like English.",
         "size": 11, "color": C_DARK},
        {"text": "", "size": 6},
        {"text": "The same syllable — four different words:",
         "size": 11, "bold": True, "color": C_NAVY},
        {"text": "mā (mother)  ·  má (hemp)  ·  mǎ (horse)  ·  mà (scold)",
         "size": 12, "bold": True, "color": C_TEAL},
        {"text": "", "size": 6},
        {"text": "The core difficulty:", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "• Tone perception is a trained skill — non-tonal language "
                 "speakers' auditory systems are not tuned to hear pitch as meaning",
         "size": 10, "color": C_DARK},
        {"text": "• Learners need high-repetition, low-stakes practice with "
                 "immediate, specific feedback",
         "size": 10, "color": C_DARK},
        {"text": "• Classroom time is limited; self-study tools rarely provide "
                 "acoustic-level guidance",
         "size": 10, "color": C_DARK},
    ])

    # Right card — The Gap
    RX = LM + PANEL_W + GAP
    rect(slide, RX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, RX, TM, PANEL_W, BAR_H, C_TEAL)
    textbox(slide, RX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "The Gap", 14, bold=True, color=C_WHITE)

    rt = TM + BAR_H + PAD * 2
    rh = PANEL_H - BAR_H - PAD * 4
    multibox(slide, RX + PAD * 2, rt, PANEL_W - PAD * 4, rh, [
        {"text": "What learners need:", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "• Instant acoustic feedback on their own voice",
         "size": 10, "color": C_DARK},
        {"text": "• Visual comparison: what they said vs. what they should say",
         "size": 10, "color": C_DARK},
        {"text": "• Plain-language guidance — not just right/wrong",
         "size": 10, "color": C_DARK},
        {"text": "• Zero installation, accessible from any browser",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 6},
        {"text": "What existing tools offer:",
         "size": 11, "bold": True, "color": C_NAVY},
        {"text": "• Most CAPT systems are desktop-only, lab-based, or "
                 "require expensive hardware",
         "size": 10, "color": C_DARK},
        {"text": "• Commercial apps (Duolingo, HelloChinese) provide "
                 "right/wrong feedback only — no pitch contour",
         "size": 10, "color": C_DARK},
        {"text": "• Research prototypes not publicly accessible",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 6},
        {"text": "→  No lightweight, browser-based tool with real-time "
                 "visual pitch feedback exists for non-tonal learners.",
         "size": 11, "bold": True, "color": C_TEAL},
    ])


# ── Slide 3: Related Work + Gap ──────────────────────────────────────────────────
def slide_related_work(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Related Work — Where This Project Fits")

    PANEL_W = 2700000
    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040
    GAP     = PAD * 2
    TOTAL_W = SW - LM * 2

    # Three equal columns
    COL_W = (TOTAL_W - GAP * 2) // 3

    cols = [
        {
            "title": "CAPT for Mandarin",
            "bar_color": C_NAVY,
            "lines": [
                {"text": "Computer-Assisted Pronunciation Training (CAPT) systems "
                         "for Mandarin have demonstrated measurable gains in "
                         "tone perception.", "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Wang et al. (2003)", "size": 10, "bold": True, "color": C_NAVY},
                {"text": "First large-scale CAPT study for Mandarin tones; "
                         "showed short-term perceptual gains with ASR feedback.",
                 "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Hincks (2003)", "size": 10, "bold": True, "color": C_NAVY},
                {"text": "Demonstrated learner benefit from visual pitch "
                         "displays in prosody training.", "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Gap: systems were lab-based, not publicly accessible, "
                         "and lacked per-tone corrective messaging.",
                 "size": 10, "italic": True, "color": C_MUTED},
            ],
        },
        {
            "title": "Automatic Tone Detection",
            "bar_color": C_MID_BLUE,
            "lines": [
                {"text": "Accurate automatic classification of Mandarin tones "
                         "from learner speech is a prerequisite for real-time "
                         "feedback.", "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Xu et al. (2004)", "size": 10, "bold": True, "color": C_NAVY},
                {"text": "F0 contour shape features (slope, height, curvature) "
                         "are the primary cue for tone classification.",
                 "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Chao (1968)", "size": 10, "bold": True, "color": C_NAVY},
                {"text": "Five-level tone letter system — the canonical "
                         "reference used to derive pitch contour targets.",
                 "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "This work: SVM RBF classifier on 27 speaker-normalized "
                         "F0 features — 91.95% speaker-independent accuracy.",
                 "size": 10, "bold": True, "color": C_TEAL},
            ],
        },
        {
            "title": "Multimodal Feedback",
            "bar_color": C_TEAL,
            "lines": [
                {"text": "Visual pitch overlays and corrective text feedback "
                         "together outperform audio-only or text-only approaches.",
                 "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Lin (2007)", "size": 10, "bold": True, "color": C_NAVY},
                {"text": "T3 sandhi (dipping → rising before another T3) — "
                         "phonological rule built into the feedback system.",
                 "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "Hincks (2003)", "size": 10, "bold": True, "color": C_NAVY},
                {"text": "Pitch displays increase learner self-monitoring and "
                         "reduce trial-and-error time.", "size": 10, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "This work: real-time pitch overlay + 12 tone-pair "
                         "corrective messages + confidence bar chart.",
                 "size": 10, "bold": True, "color": C_TEAL},
            ],
        },
    ]

    for i, col in enumerate(cols):
        lx = LM + i * (COL_W + GAP)
        rect(slide, lx, TM, COL_W, PANEL_H, C_CARD)
        rect(slide, lx, TM, COL_W, BAR_H, col["bar_color"])
        textbox(slide, lx + PAD, TM + PAD, COL_W - PAD * 2, BAR_H - PAD,
                col["title"], 13, bold=True, color=C_WHITE)
        ct = TM + BAR_H + PAD * 2
        ch = PANEL_H - BAR_H - PAD * 4
        multibox(slide, lx + PAD * 2, ct, COL_W - PAD * 4, ch, col["lines"])


# ── Slide 4: The Tool ────────────────────────────────────────────────────────────
def slide_tool(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Mandarin Tone Coach — What We Built")

    PANEL_W = 3800000
    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040

    # Left panel — pipeline
    LX = LM
    rect(slide, LX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, LX, TM, PANEL_W, BAR_H, C_NAVY)
    textbox(slide, LX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "The Pipeline", 14, bold=True, color=C_WHITE)

    ct = TM + BAR_H + PAD * 2
    ch = PANEL_H - BAR_H - PAD * 4
    multibox(slide, LX + PAD * 2, ct, PANEL_W - PAD * 4, ch, [
        {"text": "Browser-based · zero install · works on any device",
         "size": 11, "bold": True, "color": C_TEAL},
        {"text": "", "size": 5},
        {"text": "① Record", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "User records a Mandarin syllable or word in the browser",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "② Extract F0", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "Parselmouth (Praat) extracts pitch contour (75–500 Hz); "
                 "normalized to semitones relative to speaker mean",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "③ Classify", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "SVM RBF classifier · 27 features · 91.95% accuracy · "
                 "3 post-processing rules for T2/T3 edge cases",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "④ Visualize", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "Learner pitch overlay (red) vs. canonical reference (blue dashed) "
                 "+ confidence bar chart",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 4},
        {"text": "⑤ Feedback", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "Plain-language corrective message for every tone-pair combination; "
                 "T3 sandhi accepted",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 6},
        {"text": "Two modes: monosyllabic (40 words) · disyllabic (20 words)",
         "size": 11, "bold": True, "color": C_NAVY},
    ])

    # Right panel — app screenshot
    IMG_X = LM + PANEL_W + PAD * 2
    IMG_W = SW - IMG_X - LM
    IMG_H = PANEL_H

    img_path = RESULTS_DIR / "slide_01.png"
    if img_path.exists():
        image(slide, img_path, IMG_X, TM, IMG_W, IMG_H)
    else:
        rect(slide, IMG_X, TM, IMG_W, IMG_H, C_DARK2)
        textbox(slide, IMG_X + PAD * 2, TM + IMG_H // 2 - 200000,
                IMG_W - PAD * 4, 400000,
                "mandarin-tone-coach.streamlit.app",
                14, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)


# ── Slide 5: Demo ────────────────────────────────────────────────────────────────
def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SW, SH, C_BG_DARK)
    rect(slide, 0, 0, 164592, SH, C_TEAL)

    textbox(slide, 457200, 1200000, SW - 457200 - LM, 700000,
            "Live Demo", 48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    textbox(slide, 457200, 2050000, SW - 457200 - LM, 400000,
            "mandarin-tone-coach.streamlit.app",
            22, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

    rect(slide, SW // 2 - 1828800, 2600000, 3657600, 36576, C_MID_BLUE)

    multibox(slide, 457200, 2750000, SW - 457200 - LM, 600000, [
        {"text": "Switching to browser now  →  practice app · perception test",
         "size": 14, "color": C_LIGHT, "align": PP_ALIGN.CENTER},
    ])


# ── Slide 6: Study ───────────────────────────────────────────────────────────────
def slide_study(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "The Study — One Session, Three Steps")

    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040
    TOTAL_W = SW - LM * 2
    GAP     = PAD * 2
    COL_W   = (TOTAL_W - GAP * 2) // 3

    steps = [
        {
            "num": "①", "title": "Pre-test  (Form A)",
            "bar_color": C_NAVY,
            "lines": [
                {"text": "12 audio clips", "size": 13, "bold": True, "color": C_NAVY},
                {"text": "Select the Mandarin tone you hear (T1–T4)",
                 "size": 11, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "• 3 items per tone", "size": 10, "color": C_DARK},
                {"text": "• Female Voice 1 + Male Voice 1", "size": 10, "color": C_DARK},
                {"text": "• No feedback given", "size": 10, "color": C_DARK},
                {"text": "• ~2–3 minutes", "size": 10, "color": C_DARK},
                {"text": "", "size": 6},
                {"text": "n = 19 participants completed Form A",
                 "size": 11, "bold": True, "color": C_TEAL},
            ],
        },
        {
            "num": "②", "title": "Practice the App",
            "bar_color": C_BLUE,
            "lines": [
                {"text": "5–10 minutes", "size": 13, "bold": True, "color": C_NAVY},
                {"text": "Mandarin Tone Coach — free practice",
                 "size": 11, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "• Monosyllabic & disyllabic words",
                 "size": 10, "color": C_DARK},
                {"text": "• Real-time pitch contour feedback",
                 "size": 10, "color": C_DARK},
                {"text": "• Corrective text guidance", "size": 10, "color": C_DARK},
                {"text": "• Self-paced, asynchronous", "size": 10, "color": C_DARK},
                {"text": "", "size": 6},
                {"text": "16 participants also completed TAM survey",
                 "size": 11, "bold": True, "color": C_TEAL},
            ],
        },
        {
            "num": "③", "title": "Post-test  (Form B)",
            "bar_color": C_GREEN,
            "lines": [
                {"text": "12 new audio clips", "size": 13, "bold": True, "color": C_NAVY},
                {"text": "Same format, new speakers",
                 "size": 11, "color": C_DARK},
                {"text": "", "size": 5},
                {"text": "• Female Voice 2 + Male Voice 2",
                 "size": 10, "color": C_DARK},
                {"text": "• Counterbalanced to reduce item-overlap confound",
                 "size": 10, "color": C_DARK},
                {"text": "• Participant ID links Form A → Form B",
                 "size": 10, "color": C_DARK},
                {"text": "", "size": 6},
                {"text": "n = 13 paired participants (pre + post)",
                 "size": 11, "bold": True, "color": C_TEAL},
            ],
        },
    ]

    for i, step in enumerate(steps):
        lx = LM + i * (COL_W + GAP)
        rect(slide, lx, TM, COL_W, PANEL_H, C_CARD)
        rect(slide, lx, TM, COL_W, BAR_H, step["bar_color"])
        textbox(slide, lx + PAD, TM + PAD, COL_W - PAD * 2, BAR_H - PAD,
                f"{step['num']}  {step['title']}", 13, bold=True, color=C_WHITE)
        ct = TM + BAR_H + PAD * 2
        ch = PANEL_H - BAR_H - PAD * 4
        multibox(slide, lx + PAD * 2, ct, COL_W - PAD * 4, ch, step["lines"])


# ── Slide 7: Results ─────────────────────────────────────────────────────────────
def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Results — What We Found")

    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040
    GAP     = PAD * 2
    TOTAL_W = SW - LM * 2

    # Left: big stat callout panel (dark)
    STAT_W = 2800000
    rect(slide, LM, TM, STAT_W, PANEL_H, C_NAVY)

    multibox(slide, LM + PAD * 2, TM + PAD * 3, STAT_W - PAD * 4, PANEL_H - PAD * 6, [
        {"text": "Tone Perception", "size": 13, "bold": True, "color": C_LIGHT,
         "align": PP_ALIGN.CENTER},
        {"text": "+10.3 pp", "size": 38, "bold": True, "color": C_TEAL,
         "align": PP_ALIGN.CENTER},
        {"text": "mean accuracy gain", "size": 11, "color": C_LIGHT,
         "align": PP_ALIGN.CENTER},
        {"text": "53.8% → 64.1%", "size": 13, "bold": True, "color": C_WHITE,
         "align": PP_ALIGN.CENTER},
        {"text": "n = 13 paired participants", "size": 10, "color": C_MUTED,
         "align": PP_ALIGN.CENTER},
        {"text": "p = 0.087  (positive trend)", "size": 10, "italic": True,
         "color": C_MUTED, "align": PP_ALIGN.CENTER},
        {"text": "", "size": 8},
        {"text": "Usability (TAM)", "size": 13, "bold": True, "color": C_LIGHT,
         "align": PP_ALIGN.CENTER},
        {"text": "3.96 / 5.0", "size": 38, "bold": True, "color": C_TEAL,
         "align": PP_ALIGN.CENTER},
        {"text": "overall mean  ·  n = 16", "size": 11, "color": C_LIGHT,
         "align": PP_ALIGN.CENTER},
        {"text": "4.19 — Would recommend", "size": 11, "bold": True,
         "color": C_WHITE, "align": PP_ALIGN.CENTER},
    ])

    # Right: per-tone breakdown + survey highlights
    RX = LM + STAT_W + GAP
    RIGHT_W = TOTAL_W - STAT_W - GAP

    rect(slide, RX, TM, RIGHT_W, PANEL_H, C_CARD)
    rect(slide, RX, TM, RIGHT_W, BAR_H, C_TEAL)
    textbox(slide, RX + PAD, TM + PAD, RIGHT_W - PAD * 2, BAR_H - PAD,
            "Per-Tone Accuracy Gain", 14, bold=True, color=C_WHITE)

    ct = TM + BAR_H + PAD * 2
    ch = PANEL_H - BAR_H - PAD * 4
    multibox(slide, RX + PAD * 2, ct, RIGHT_W - PAD * 4, ch, [
        {"text": "T3  Dipping   46.2% → 66.7%   +20.5 pp  ▲ largest",
         "size": 12, "bold": True, "color": C_TEAL},
        {"text": "Consistent with literature — dipping tone hardest for non-tonal learners  (Lin, 2007)",
         "size": 10, "italic": True, "color": C_MUTED},
        {"text": "", "size": 5},
        {"text": "T1  Flat       53.8% → 69.2%   +15.4 pp  ▲",
         "size": 12, "color": C_NAVY},
        {"text": "", "size": 3},
        {"text": "T2  Rising    53.8% → 66.7%   +12.8 pp  ▲",
         "size": 12, "color": C_NAVY},
        {"text": "", "size": 3},
        {"text": "T4  Falling   61.5% → 53.8%    −7.7 pp  ▼",
         "size": 12, "color": C_RED},
        {"text": "3 items/tone/participant — likely sampling variability",
         "size": 10, "italic": True, "color": C_MUTED},
        {"text": "", "size": 8},
        {"text": "Survey Highlights", "size": 12, "bold": True, "color": C_NAVY},
        {"text": "• Would recommend to other learners:  4.19 / 5",
         "size": 11, "color": C_DARK},
        {"text": "• Overall useful for learning tones:   4.12 / 5",
         "size": 11, "color": C_DARK},
        {"text": "• Pitch contour helped identify errors:  3.88 / 5",
         "size": 11, "color": C_DARK},
        {"text": "• Note: 3 native/heritage speakers in sample may depress ratings",
         "size": 10, "italic": True, "color": C_MUTED},
    ])


# ── Slide 8: Conclusion + Future Work ────────────────────────────────────────────
def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Conclusion & Future Work")

    PANEL_W = 4115520
    PANEL_H = SH - TM - PAD * 2
    BAR_H   = 320040
    GAP     = PAD * 2

    # Left: Conclusion / what worked
    LX = LM
    rect(slide, LX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, LX, TM, PANEL_W, BAR_H, C_NAVY)
    textbox(slide, LX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "What We Showed", 14, bold=True, color=C_WHITE)

    ct = TM + BAR_H + PAD * 2
    ch = PANEL_H - BAR_H - PAD * 4
    multibox(slide, LX + PAD * 2, ct, PANEL_W - PAD * 4, ch, [
        {"text": "A lightweight, browser-based CAPT tool can produce "
                 "meaningful perceptual gains with minimal instructional overhead.",
         "size": 11, "bold": True, "color": C_NAVY},
        {"text": "", "size": 6},
        {"text": "✓  +10.3 pp improvement in one session (n = 13)",
         "size": 11, "color": C_DARK},
        {"text": "✓  T3 (dipping) improved most — the hardest tone at baseline",
         "size": 11, "color": C_DARK},
        {"text": "✓  91.95% SVM classifier — competitive with prior work",
         "size": 11, "color": C_DARK},
        {"text": "✓  TAM 3.96 / 5  —  participants would recommend the app",
         "size": 11, "color": C_DARK},
        {"text": "✓  T2 / T3 confusion in learners mirrors classifier errors — "
                 "a useful diagnostic signal",
         "size": 11, "color": C_DARK},
        {"text": "", "size": 6},
        {"text": "Limitations", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "• n = 13 paired — underpowered; ~28 needed for 80% power",
         "size": 10, "color": C_DARK},
        {"text": "• No control group — can't rule out test-retest effects",
         "size": 10, "color": C_DARK},
        {"text": "• Single session — no long-term retention data",
         "size": 10, "color": C_DARK},
        {"text": "• Convenience sample (same class) — low diversity",
         "size": 10, "color": C_DARK},
    ])

    # Right: Future work
    RX = LM + PANEL_W + GAP
    rect(slide, RX, TM, PANEL_W, PANEL_H, C_CARD)
    rect(slide, RX, TM, PANEL_W, BAR_H, C_TEAL)
    textbox(slide, RX + PAD, TM + PAD, PANEL_W - PAD * 2, BAR_H - PAD,
            "What's Next", 14, bold=True, color=C_WHITE)

    rt = TM + BAR_H + PAD * 2
    rh = PANEL_H - BAR_H - PAD * 4
    multibox(slide, RX + PAD * 2, rt, PANEL_W - PAD * 4, rh, [
        {"text": "Study Design", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "• Larger, more diverse participant pool",
         "size": 10, "color": C_DARK},
        {"text": "• Matched control group (practice without feedback)",
         "size": 10, "color": C_DARK},
        {"text": "• Delayed post-test — measure retention after 1 week",
         "size": 10, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "System Improvements", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "• Reference audio playback before recording",
         "size": 10, "color": C_DARK},
        {"text": "• More specific recording quality diagnostics",
         "size": 10, "color": C_DARK},
        {"text": "• Adaptive difficulty — target weakest tones first",
         "size": 10, "color": C_DARK},
        {"text": "• Speaker normalization for atypical vocal profiles",
         "size": 10, "color": C_DARK},
        {"text": "• Expanded word list", "size": 10, "color": C_DARK},
        {"text": "", "size": 5},
        {"text": "Broader Question", "size": 11, "bold": True, "color": C_NAVY},
        {"text": "Does visual pitch feedback produce durable perceptual gains, "
                 "or only short-term improvement? That's the study worth running next.",
         "size": 10, "italic": True, "color": C_TEAL},
    ])


# ── Slide 9: Thank You ────────────────────────────────────────────────────────────
def slide_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SW, SH, C_BG_DARK)
    rect(slide, 0, 0, 164592, SH, C_TEAL)

    textbox(slide, 457200, 900000, SW - 457200 - LM, 600000,
            "Thank You", 44, bold=True, color=C_WHITE)

    textbox(slide, 457200, 1600000, SW - 457200 - LM, 400000,
            "Read the paper for the full story — methodology, classifier details, and complete results.",
            16, color=C_LIGHT)

    rect(slide, 457200, 2150000, 3657600, 36576, C_MID_BLUE)

    multibox(slide, 457200, 2280000, SW - 457200 - LM, 1200000, [
        {"text": "Try the app:", "size": 13, "bold": True, "color": C_TEAL},
        {"text": "mandarin-tone-coach.streamlit.app", "size": 13, "color": C_WHITE,
         "space_before": 2},
        {"text": "", "size": 8},
        {"text": "Bharath Jagadish  ·  bjagadish3@gatech.edu",
         "size": 12, "color": C_LIGHT, "space_before": 2},
        {"text": "CS6460 Educational Technology  ·  Georgia Tech",
         "size": 12, "color": C_LIGHT, "space_before": 2},
        {"text": "", "size": 10},
        {"text": "Key References", "size": 11, "bold": True, "color": C_TEAL},
        {"text": "Chao (1968) · Wang et al. (2003) · Hincks (2003) · Xu et al. (2004) · Lin (2007)",
         "size": 10, "color": C_MUTED, "space_before": 2},
    ])


# ── Main ─────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = E(SW)
    prs.slide_height = E(SH)

    slide_title(prs)
    slide_problem(prs)
    slide_related_work(prs)
    slide_tool(prs)
    slide_demo(prs)
    slide_study(prs)
    slide_results(prs)
    slide_conclusion(prs)
    slide_thankyou(prs)

    prs.save(OUT_FILE)
    print(f"Saved: {OUT_FILE}  ({len(prs.slides)} slides)")
